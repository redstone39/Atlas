#!/usr/bin/env python3
"""Generate the nine real files used by the Atlas multiformat acceptance smoke."""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
from pathlib import Path
import subprocess
import tempfile

from docx import Document
from docx.shared import Inches
from openpyxl import Workbook
from openpyxl.drawing.image import Image as SpreadsheetImage
from openpyxl.worksheet.table import Table, TableStyleInfo
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches as SlideInches


FIXTURE_TEXT = {
    "pdf": "ORION controller target is 47 ohm.",
    "docx": "ORION tolerance is plus or minus 5 percent.",
    "pptx": "ORION board family is Atlas-One.",
    "xlsx": "ORION current limit is 2.4 ampere.",
    "txt": "ORION assembly note requires a blue inspection tag.",
    "csv": "ORION component R47 is approved.",
}


def _diagram(path: Path) -> None:
    image = Image.new("RGB", (640, 320), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((40, 105, 210, 215), outline="navy", width=5)
    draw.rectangle((430, 55, 600, 145), outline="darkgreen", width=5)
    draw.rectangle((430, 185, 600, 275), outline="darkgreen", width=5)
    draw.line((210, 160, 430, 100), fill="black", width=5)
    draw.line((210, 160, 430, 230), fill="black", width=5)
    draw.text((70, 145), "ORION CTRL", fill="navy")
    draw.text((470, 90), "SENSOR A", fill="darkgreen")
    draw.text((470, 220), "SENSOR B", fill="darkgreen")
    image.save(path, format="PNG", optimize=False)


def _pdf_source(path: Path) -> None:
    document = Document()
    document.add_heading("ORION controller target", level=1)
    document.add_paragraph(FIXTURE_TEXT["pdf"])
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Parameter"
    table.rows[0].cells[1].text = "Value"
    cells = table.add_row().cells
    cells[0].text = "Target resistance"
    cells[1].text = "47 ohm"
    document.save(path)


def _docx(path: Path, diagram: Path) -> None:
    document = Document()
    document.add_heading("ORION qualification", level=1)
    document.add_paragraph(FIXTURE_TEXT["docx"])
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Parameter"
    table.rows[0].cells[1].text = "Value"
    for name, value in (
        ("Qualification tolerance", "plus or minus 5 percent"),
        ("Inspection", "blue tag"),
    ):
        cells = table.add_row().cells
        cells[0].text = name
        cells[1].text = value
    document.add_picture(str(diagram), width=Inches(4.6))
    document.add_page_break()
    document.add_heading("ORION qualification evidence", level=1)
    document.add_paragraph("The qualification record was reviewed by QA.")
    document.add_paragraph("The repeated anchor belongs to this page only.")
    document.add_picture(str(diagram), width=Inches(3.8))
    document.add_paragraph("The repeated anchor belongs to this page only.")
    document.save(path)


def _pptx(path: Path, diagram: Path) -> None:
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "ORION board family"
    textbox = first.shapes.add_textbox(
        SlideInches(0.7), SlideInches(1.3), SlideInches(5.5), SlideInches(0.7)
    )
    textbox.text_frame.text = FIXTURE_TEXT["pptx"]
    first.shapes.add_picture(
        str(diagram), SlideInches(0.7), SlideInches(2.1), width=SlideInches(5.0)
    )

    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "ORION sensor allocation"
    chart_data = ChartData()
    chart_data.categories = ["Sensor A", "Sensor B"]
    chart_data.add_series("Channels", (2, 3))
    second.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        SlideInches(0.7),
        SlideInches(1.5),
        SlideInches(5.4),
        SlideInches(3.5),
        chart_data,
    )
    note = second.shapes.add_textbox(
        SlideInches(6.3), SlideInches(1.7), SlideInches(3.0), SlideInches(1.5)
    )
    note.text_frame.text = "Controller topology diagram"
    second.shapes.add_picture(
        str(diagram), SlideInches(6.3), SlideInches(3.1), width=SlideInches(3.0)
    )
    presentation.save(path)


def _xlsx(path: Path, diagram: Path) -> None:
    workbook = Workbook()
    specs = workbook.active
    specs.title = "Specs"
    rows = [
        ("Parameter", "Value"),
        ("Current limit", "2.4 ampere"),
        ("Inspection tag", "blue"),
        ("Board family", "Atlas-One"),
    ]
    for row in rows:
        specs.append(row)
    table = Table(displayName="OrionSpecs", ref="A1:B4")
    table.tableStyleInfo = TableStyleInfo(
        name="TableStyleMedium2", showRowStripes=True
    )
    specs.add_table(table)
    specs["A6"] = FIXTURE_TEXT["xlsx"]
    specs.add_image(SpreadsheetImage(str(diagram)), "D2")
    specs.print_area = "A1:H28"

    limits = workbook.create_sheet("Limits")
    limits.append(("Limit", "Value"))
    limits.append(("Sensor A channels", "2"))
    limits.append(("Sensor B channels", "3"))
    limits.add_image(SpreadsheetImage(str(diagram)), "D4")
    limits.print_area = "A1:H28"
    workbook.save(path)


def _convert(soffice: str, source: Path, target_filter: str, output: Path) -> None:
    generated = output.parent / f"{source.stem}.{target_filter.split(':', 1)[0]}"
    with tempfile.TemporaryDirectory(prefix="atlas-fixture-lo-") as profile:
        command = [
            soffice,
            "--headless",
            f"-env:UserInstallation={Path(profile).resolve().as_uri()}",
            "--convert-to",
            target_filter,
            "--outdir",
            str(output.parent),
            str(source),
        ]
        completed = subprocess.run(
            command, check=False, capture_output=True, text=True, timeout=120
        )
    if completed.returncode == 0 and generated.is_file() and generated != output:
        generated.replace(output)
    if completed.returncode != 0 or not output.is_file() or output.stat().st_size == 0:
        raise RuntimeError(
            f"LibreOffice conversion failed for {source.name}: "
            f"{completed.stdout[-500:]} {completed.stderr[-500:]}"
        )


def _digest(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--soffice", default="soffice")
    args = parser.parse_args()
    output = args.output.resolve()
    output.mkdir(parents=True, exist_ok=True)

    diagram = output / "orion-diagram.png"
    pdf_source = output / "orion-target-source.docx"
    docx = output / "orion-qualification.docx"
    pptx = output / "orion-board.pptx"
    xlsx = output / "orion-current.xlsx"
    _diagram(diagram)
    _pdf_source(pdf_source)
    _docx(docx, diagram)
    _pptx(pptx, diagram)
    _xlsx(xlsx, diagram)

    (output / "orion-notes.txt").write_text(
        FIXTURE_TEXT["txt"] + "\n", encoding="utf-8"
    )
    with (output / "orion-components.csv").open(
        "w", encoding="utf-8", newline=""
    ) as stream:
        writer = csv.writer(stream)
        writer.writerow(("component", "status", "note"))
        writer.writerow(("R47", "approved", FIXTURE_TEXT["csv"]))
        writer.writerow(("CTRL", "qualified", "Atlas-One controller"))

    _convert(
        args.soffice,
        pdf_source,
        "pdf:writer_pdf_Export",
        output / "orion-target.pdf",
    )
    _convert(
        args.soffice,
        docx,
        "doc:MS Word 97",
        output / "orion-qualification-legacy.doc",
    )
    _convert(
        args.soffice,
        pptx,
        "ppt:MS PowerPoint 97",
        output / "orion-board-legacy.ppt",
    )
    _convert(
        args.soffice,
        xlsx,
        "xls:MS Excel 97",
        output / "orion-current-legacy.xls",
    )

    expected = {
        "pdf": "orion-target.pdf",
        "docx": docx.name,
        "pptx": pptx.name,
        "xlsx": xlsx.name,
        "txt": "orion-notes.txt",
        "csv": "orion-components.csv",
        "doc": "orion-qualification-legacy.doc",
        "ppt": "orion-board-legacy.ppt",
        "xls": "orion-current-legacy.xls",
    }
    manifest = {
        "schema_version": "atlas-multiformat-fixtures-v1",
        "files": {
            document_format: {
                "filename": filename,
                "sha256": _digest(output / filename),
                "byte_length": (output / filename).stat().st_size,
            }
            for document_format, filename in expected.items()
        },
        "complementary_query": (
            "For ORION, what values are listed for Target resistance and "
            "Qualification tolerance?"
        ),
        "expected_answer_terms": ["47 ohm", "plus or minus 5 percent"],
        "visual_query": "What does ORION VISUAL BUS connect?",
    }
    (output / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )
    diagram.unlink()
    pdf_source.unlink()
    print(json.dumps(manifest, ensure_ascii=False, sort_keys=True))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
