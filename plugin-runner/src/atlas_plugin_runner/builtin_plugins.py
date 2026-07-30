from __future__ import annotations

import csv
from io import BytesIO
import json
import os
import re

from atlas_processing_sdk import CandidateDraft, SourceRegionDraft


class PypdfPlugin:
    async def parse(self, request, context):
        pages = await context.artifact_broker.parsed_pdf_pages(
            request.artifact_ref, request.unit_start, request.unit_end
        )
        for page_number, text_ref in pages:
            yield SourceRegionDraft(
                source_region_identity=f"page:{page_number}",
                region_kind="page",
                content_kind_hint="text" if text_ref else "unknown",
                element_kind_hint="page",
                locator_draft={"selector_kind": "page_region", "page_number": page_number},
                normalized_text_ref=text_ref,
            )
            if not text_ref:
                continue
            page_text = await context.artifact_broker.read_text(text_ref)
            paragraphs = [
                value.strip()
                for value in re.split(r"\n\s*\n", page_text.strip())
                if value.strip()
            ]
            for ordinal, paragraph in enumerate(paragraphs, start=1):
                paragraph_ref = context.artifact_broker.put_text(paragraph)
                yield SourceRegionDraft(
                    source_region_identity=f"page:{page_number}:paragraph:{ordinal}",
                    parent_region_identity=f"page:{page_number}",
                    region_kind="paragraph",
                    content_kind_hint="text",
                    element_kind_hint="paragraph",
                    locator_draft={
                        "selector_kind": "normalized_text_span",
                        "page_number": page_number,
                        "ordinal": ordinal,
                    },
                    normalized_text_ref=paragraph_ref,
                )


class InlineTextPlugin:
    async def parse(self, request, context):
        text = (await context.artifact_broker.read_bytes(request.artifact_ref)).decode("utf-8-sig").strip()
        if not text:
            return
        paragraph_start = 0
        for ordinal, match in enumerate(
            re.finditer(r"(?:^|\n\s*\n)(.*?)(?=\n\s*\n|\Z)", text, re.S),
            start=1,
        ):
            paragraph = match.group(1).strip()
            if not paragraph:
                continue
            paragraph_start = text[: match.start(1)].count("\n") + 1
            paragraph_end = paragraph_start + paragraph.count("\n")
            ref = context.artifact_broker.put_text(paragraph)
            yield SourceRegionDraft(
                source_region_identity=f"paragraph:{ordinal}",
                region_kind="paragraph",
                content_kind_hint="text",
                element_kind_hint="paragraph",
                locator_draft={
                    "selector_kind": "text_line_span",
                    "ordinal": ordinal,
                    "line_start": paragraph_start,
                    "line_end": paragraph_end,
                },
                normalized_text_ref=ref,
            )


class CsvPlugin:
    async def parse(self, request, context):
        source = await context.artifact_broker.read_bytes(request.artifact_ref)
        text = source.decode("utf-8-sig")
        rows = list(csv.reader(text.splitlines()))
        if not rows:
            return
        width = max((len(row) for row in rows), default=0)
        header = rows[0]
        for row_number, row in enumerate(rows, start=1):
            normalized = "\t".join(value.strip() for value in row).strip()
            if not normalized:
                continue
            structured_ref = context.artifact_broker.put_text(json.dumps(
                {
                    "kind": "csv_row",
                    "row_number": row_number,
                    "header": header,
                    "cells": [
                        {"column_index": index, "value": value}
                        for index, value in enumerate(row, start=1)
                    ],
                },
                ensure_ascii=False,
                sort_keys=True,
                separators=(",", ":"),
            ))
            yield SourceRegionDraft(
                source_region_identity=f"row:{row_number}",
                region_kind="table",
                content_kind_hint="table",
                element_kind_hint="header" if row_number == 1 else "row",
                locator_draft={
                    "selector_kind": "csv_row",
                    "row_number": row_number,
                    "column_count": len(row),
                    "table_width": width,
                },
                normalized_text_ref=context.artifact_broker.put_text(normalized),
                structured_content_ref=structured_ref,
            )


class DocxPlugin:
    async def parse(self, request, context):
        from docx import Document

        document = Document(BytesIO(
            await context.artifact_broker.read_bytes(request.artifact_ref)
        ))
        for paragraph_index, paragraph in enumerate(document.paragraphs, start=1):
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = paragraph.style.name if paragraph.style is not None else ""
            num_pr = paragraph._p.pPr.numPr if paragraph._p.pPr is not None else None
            paragraph_kind = (
                "heading" if style_name.casefold().startswith("heading")
                else "list_item" if num_pr is not None
                else "paragraph"
            )
            yield SourceRegionDraft(
                source_region_identity=f"paragraph:{paragraph_index}",
                region_kind="paragraph",
                content_kind_hint="text",
                element_kind_hint=paragraph_kind,
                locator_draft={
                    "selector_kind": "word_paragraph",
                    "paragraph_index": paragraph_index,
                    "paragraph_kind": paragraph_kind,
                    "style_name": style_name,
                },
                normalized_text_ref=context.artifact_broker.put_text(text),
            )
        for table_index, table in enumerate(document.tables, start=1):
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            text = "\n".join("\t".join(row) for row in rows).strip()
            if not text:
                continue
            structured = {
                "kind": "word_table",
                "table_index": table_index,
                "rows": [
                    [
                        {"row_index": r, "column_index": c, "value": value}
                        for c, value in enumerate(row, start=1)
                    ]
                    for r, row in enumerate(rows, start=1)
                ],
            }
            yield SourceRegionDraft(
                source_region_identity=f"table:{table_index}",
                region_kind="table",
                content_kind_hint="table",
                element_kind_hint="table",
                locator_draft={
                    "selector_kind": "word_table",
                    "table_index": table_index,
                    "row_count": len(rows),
                    "column_count": max((len(row) for row in rows), default=0),
                },
                normalized_text_ref=context.artifact_broker.put_text(text),
                structured_content_ref=context.artifact_broker.put_text(json.dumps(
                    structured, ensure_ascii=False, sort_keys=True, separators=(",", ":")
                )),
            )
        body_blocks: list[tuple[str, list[str]]] = []
        for child in document.element.body.iterchildren():
            text = " ".join(
                value.text.strip()
                for value in child.iter()
                if value.tag.endswith("}t")
                and isinstance(value.text, str)
                and value.text.strip()
            )
            relationship_ids = [
                value
                for value_node in child.iter()
                for attribute, value in value_node.attrib.items()
                if attribute.endswith("}embed")
                and isinstance(value, str)
                and value
            ]
            body_blocks.append((text, relationship_ids))

        image_index = 0
        for block_index, (_text, relationship_ids) in enumerate(body_blocks):
            previous_anchor = next(
                (
                    text for text, _ in reversed(body_blocks[:block_index])
                    if text
                ),
                None,
            )
            following_anchor = next(
                (
                    text for text, _ in body_blocks[block_index + 1 :]
                    if text
                ),
                None,
            )
            alignment_anchors = [
                value for value in (previous_anchor, following_anchor) if value
            ]
            for relationship_id in relationship_ids:
                part = document.part.related_parts.get(relationship_id)
                content_type = getattr(part, "content_type", "")
                blob = getattr(part, "blob", None)
                if not content_type.startswith("image/") or not isinstance(blob, bytes):
                    continue
                image_index += 1
                yield SourceRegionDraft(
                    source_region_identity=f"image:{image_index}",
                    region_kind="image_region",
                    content_kind_hint="image",
                    element_kind_hint="image",
                    locator_draft={
                        "selector_kind": "word_image",
                        "image_index": image_index,
                        "relationship_id": relationship_id,
                        "part_name": str(getattr(part, "partname", "")).lstrip("/"),
                        "alignment_anchors": alignment_anchors,
                    },
                    native_artifact_ref=context.artifact_broker.put_bytes(blob),
                )


class PptxPlugin:
    async def parse(self, request, context):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        presentation = Presentation(BytesIO(
            await context.artifact_broker.read_bytes(request.artifact_ref)
        ))
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text.strip():
                    slide_text.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    slide_text.extend(
                        cell.text.strip()
                        for row in shape.table.rows for cell in row.cells
                        if cell.text.strip()
                    )
            slide_ref = (
                context.artifact_broker.put_text("\n".join(slide_text))
                if slide_text else None
            )
            yield SourceRegionDraft(
                source_region_identity=f"slide:{slide_number}",
                region_kind="slide",
                content_kind_hint="text" if slide_ref else "unknown",
                element_kind_hint="slide",
                locator_draft={
                    "selector_kind": "powerpoint_slide",
                    "slide_number": slide_number,
                    "slide_width": int(presentation.slide_width),
                    "slide_height": int(presentation.slide_height),
                },
                normalized_text_ref=slide_ref,
            )
            for shape_index, shape in enumerate(slide.shapes, start=1):
                geometry = {
                    "left": int(shape.left), "top": int(shape.top),
                    "width": int(shape.width), "height": int(shape.height),
                }
                base_locator = {
                    "slide_number": slide_number,
                    "slide_width": int(presentation.slide_width),
                    "slide_height": int(presentation.slide_height),
                    "shape_index": shape_index,
                    "shape_name": shape.name,
                    **geometry,
                }
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    yield SourceRegionDraft(
                        source_region_identity=f"slide:{slide_number}:image:{shape_index}",
                        parent_region_identity=f"slide:{slide_number}",
                        region_kind="image_region",
                        content_kind_hint="image",
                        element_kind_hint="image",
                        locator_draft={"selector_kind": "powerpoint_image", **base_locator},
                        native_artifact_ref=context.artifact_broker.put_bytes(shape.image.blob),
                    )
                    continue
                if getattr(shape, "has_table", False):
                    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    text = "\n".join("\t".join(row) for row in rows).strip()
                    if text:
                        yield SourceRegionDraft(
                            source_region_identity=f"slide:{slide_number}:table:{shape_index}",
                            parent_region_identity=f"slide:{slide_number}",
                            region_kind="table",
                            content_kind_hint="table",
                            element_kind_hint="table",
                            locator_draft={"selector_kind": "powerpoint_table", **base_locator},
                            normalized_text_ref=context.artifact_broker.put_text(text),
                            structured_content_ref=context.artifact_broker.put_text(json.dumps(
                                {"kind": "powerpoint_table", "rows": rows},
                                ensure_ascii=False, sort_keys=True, separators=(",", ":"),
                            )),
                        )
                    continue
                if getattr(shape, "has_chart", False):
                    chart = shape.chart
                    values: list[str] = []
                    if chart.has_title and chart.chart_title.has_text_frame:
                        values.append(chart.chart_title.text_frame.text.strip())
                    for series in chart.series:
                        values.append(str(series.name))
                        values.extend(str(value) for value in series.values)
                    text = " ".join(value for value in values if value).strip()
                    if text:
                        yield SourceRegionDraft(
                            source_region_identity=f"slide:{slide_number}:chart:{shape_index}",
                            parent_region_identity=f"slide:{slide_number}",
                            region_kind="figure",
                            content_kind_hint="figure",
                            element_kind_hint="chart",
                            locator_draft={"selector_kind": "powerpoint_chart", **base_locator},
                            normalized_text_ref=context.artifact_broker.put_text(text),
                        )
                    continue
                text = shape.text.strip() if getattr(shape, "has_text_frame", False) else ""
                if text:
                    yield SourceRegionDraft(
                        source_region_identity=f"slide:{slide_number}:shape:{shape_index}",
                        parent_region_identity=f"slide:{slide_number}",
                        region_kind="paragraph",
                        content_kind_hint="text",
                        element_kind_hint="shape_text",
                        locator_draft={"selector_kind": "powerpoint_shape", **base_locator},
                        normalized_text_ref=context.artifact_broker.put_text(text),
                    )


class XlsxPlugin:
    async def parse(self, request, context):
        from openpyxl import load_workbook

        workbook = load_workbook(
            BytesIO(await context.artifact_broker.read_bytes(request.artifact_ref)),
            data_only=False,
        )
        for sheet_index, sheet in enumerate(workbook.worksheets, start=1):
            nonempty_row_text: dict[int, str] = {}
            for row_index, row in enumerate(sheet.iter_rows(), start=1):
                cells = [
                    {"coordinate": cell.coordinate, "value": "" if cell.value is None else str(cell.value)}
                    for cell in row if cell.value is not None
                ]
                if not cells:
                    continue
                text = "\t".join(cell["value"] for cell in cells)
                if not text.strip():
                    continue
                nonempty_row_text[row_index] = text
                yield SourceRegionDraft(
                    source_region_identity=f"sheet:{sheet_index}:row:{row_index}",
                    region_kind="table",
                    content_kind_hint="table",
                    element_kind_hint="row",
                    locator_draft={
                        "selector_kind": "excel_row",
                        "sheet_index": sheet_index,
                        "sheet_name": sheet.title,
                        "row_index": row_index,
                        "cell_coordinates": [cell["coordinate"] for cell in cells],
                    },
                    normalized_text_ref=context.artifact_broker.put_text(text),
                    structured_content_ref=context.artifact_broker.put_text(json.dumps(
                        {"kind": "excel_row", "cells": cells},
                        ensure_ascii=False, sort_keys=True, separators=(",", ":"),
                    )),
                )
            for table_index, table in enumerate(sheet.tables.values(), start=1):
                yield SourceRegionDraft(
                    source_region_identity=f"sheet:{sheet_index}:table:{table_index}",
                    region_kind="table",
                    content_kind_hint="table",
                    element_kind_hint="table",
                    locator_draft={
                        "selector_kind": "excel_table",
                        "sheet_index": sheet_index,
                        "sheet_name": sheet.title,
                        "table_index": table_index,
                        "table_name": table.name,
                        "cell_range": table.ref,
                    },
                    normalized_text_ref=context.artifact_broker.put_text(
                        f"{table.name} {table.ref}"
                    ),
                )
            for image_index, image in enumerate(getattr(sheet, "_images", ()), start=1):
                anchor = getattr(image, "anchor", None)
                marker = getattr(anchor, "_from", None)
                anchor_row = int(marker.row) + 1 if marker is not None else None
                prior_rows = [
                    row_index for row_index in nonempty_row_text
                    if anchor_row is not None and row_index < anchor_row
                ]
                following_rows = [
                    row_index for row_index in nonempty_row_text
                    if anchor_row is not None and row_index > anchor_row
                ]
                alignment_anchors = [
                    nonempty_row_text[row_index]
                    for row_index in (
                        max(prior_rows) if prior_rows else None,
                        min(following_rows) if following_rows else None,
                    )
                    if row_index is not None
                ]
                yield SourceRegionDraft(
                    source_region_identity=f"sheet:{sheet_index}:image:{image_index}",
                    region_kind="image_region",
                    content_kind_hint="image",
                    element_kind_hint="image",
                    locator_draft={
                        "selector_kind": "excel_image",
                        "sheet_index": sheet_index,
                        "sheet_name": sheet.title,
                        "image_index": image_index,
                        "anchor_row": anchor_row,
                        "anchor_column": int(marker.col) + 1 if marker is not None else None,
                        "alignment_anchors": alignment_anchors,
                    },
                    native_artifact_ref=context.artifact_broker.put_bytes(image._data()),
                )


class GenericTextPlugin:
    async def process(self, request, context):
        if not request.normalized_text_ref:
            return
        await context.artifact_broker.read_text(request.normalized_text_ref)
        yield CandidateDraft(
            source_region_ids=(request.region_id,),
            channel_id="generic_text",
            output_contract_version="eir-draft-v1",
            candidate_payload_ref=request.normalized_text_ref,
            content_kind_hint=request.content_kind_hint,
            element_kind_hint=request.element_kind_hint,
        )


def _lossless_png(source: bytes) -> tuple[bytes, int, int]:
    from PIL import Image

    with Image.open(BytesIO(source)) as image:
        normalized = image.convert("RGB")
        output = BytesIO()
        normalized.save(output, format="PNG", optimize=False, compress_level=9)
        return output.getvalue(), normalized.width, normalized.height


class RapidOcrPlugin:
    """Local OCR for embedded Office images; empty OCR is a valid no-candidate result."""

    async def process(self, request, context):
        if not request.native_artifact_ref:
            return
        from rapidocr import RapidOCR

        source = await context.artifact_broker.read_bytes(request.native_artifact_ref)
        png, _width, _height = _lossless_png(source)
        # Each invocation runs in its own isolated child. Keep ONNX from
        # creating a host-sized thread pool per child so two document workers
        # can OCR concurrently without exhausting the runner process limit.
        output = RapidOCR(params={
            "EngineConfig.onnxruntime.intra_op_num_threads": 1,
            "EngineConfig.onnxruntime.inter_op_num_threads": 1,
        })(png)
        text = "\n".join(value.strip() for value in (output.txts or ()) if value.strip())
        if not text:
            return
        yield CandidateDraft(
            source_region_ids=(request.region_id,),
            channel_id="generic_text",
            output_contract_version="eir-draft-v1",
            candidate_payload_ref=context.artifact_broker.put_text(text),
            content_kind_hint="text",
            element_kind_hint="ocr_text",
            native_artifact_ref=context.artifact_broker.put_bytes(png),
            quality_flag_refs=("ocr_local_rapidocr",),
        )


def _docling_preview_region(document, item, page_number: int, region_kind: str):
    provenance = next(
        (value for value in item.prov if value.page_no == page_number),
        None,
    )
    page = document.pages.get(page_number)
    if provenance is None or page is None or provenance.bbox is None:
        return None
    bbox = provenance.bbox.to_bottom_left_origin(page.size.height)
    return {
        "page_number": page_number,
        "region_kind": region_kind,
        "source_element_id": str(item.self_ref),
        "coordinate_system": "pdf_crop_box_relative_bottom_left",
        "rectangles": [[bbox.l, bbox.b, bbox.r, bbox.t]],
        "page_width": page.size.width,
        "page_height": page.size.height,
        "geometry_version": "docling-page-region-v1",
    }


class DoclingLayoutPlugin:
    """Trusted local/OCR layout processor for best-effort citation regions."""

    async def process(self, request, context):
        from docling.datamodel.accelerator_options import (
            AcceleratorDevice,
            AcceleratorOptions,
        )
        from docling.datamodel.base_models import (
            ConversionStatus,
            DocumentStream,
            InputFormat,
        )
        from docling.datamodel.pipeline_options import (
            PdfPipelineOptions,
            RapidOcrOptions,
        )
        from docling.document_converter import DocumentConverter, PdfFormatOption
        from docling_core.types.doc import DocItemLabel, TextItem

        source = await context.artifact_broker.read_bytes(request.artifact_ref)
        options = PdfPipelineOptions(
            artifacts_path=os.environ.get("ATLAS_DOCLING_ARTIFACTS_PATH"),
            do_ocr=True,
            enable_remote_services=False,
            allow_external_plugins=False,
            accelerator_options=AcceleratorOptions(
                num_threads=4,
                device=AcceleratorDevice.CPU,
            ),
            ocr_options=RapidOcrOptions(rapidocr_params={
                "EngineConfig.onnxruntime.intra_op_num_threads": 2,
                "EngineConfig.onnxruntime.inter_op_num_threads": 1,
            }),
        )
        page_number = request.locator_draft.get("page_number")
        artifact_page_number = getattr(request, "plugin_config", {}).get(
            "artifact_page_number"
        )
        if not isinstance(artifact_page_number, int):
            artifact_page_number = page_number
        page_range = (
            (artifact_page_number, artifact_page_number)
            if isinstance(artifact_page_number, int)
            else None
        )
        conversion = DocumentConverter(format_options={
            InputFormat.PDF: PdfFormatOption(pipeline_options=options)
        }).convert(
            DocumentStream(name="document.pdf", stream=BytesIO(source)),
            **({"page_range": page_range} if page_range is not None else {}),
        )
        if conversion.status is not ConversionStatus.SUCCESS:
            raise RuntimeError("docling_conversion_incomplete")
        document = conversion.document

        for item, _level in document.iterate_items():
            if (
                not isinstance(item, TextItem)
                or item.label == DocItemLabel.CAPTION
                or not item.text.strip()
            ):
                continue
            item_page = next((value.page_no for value in item.prov), None)
            if not isinstance(item_page, int) or (
                isinstance(artifact_page_number, int)
                and item_page != artifact_page_number
            ):
                continue
            preview = _docling_preview_region(
                document, item, item_page, "paragraph"
            )
            if preview is not None and isinstance(page_number, int):
                preview["page_number"] = page_number
            yield CandidateDraft(
                source_region_ids=(request.region_id,),
                channel_id="generic_text",
                output_contract_version="eir-draft-v1",
                candidate_payload_ref=context.artifact_broker.put_text(
                    item.text.strip()
                ),
                content_kind_hint="text",
                element_kind_hint="paragraph",
                preview_region=preview,
                quality_flag_refs=(
                    () if preview else ("pdf_preview_region_missing",)
                ),
            )

        for table_index, table in enumerate(document.tables, start=1):
            table_page = next((value.page_no for value in table.prov), None)
            if not isinstance(table_page, int) or (
                isinstance(artifact_page_number, int)
                and table_page != artifact_page_number
            ):
                continue
            page = document.pages.get(table_page)
            rows: list[list[dict[str, str]]] = [
                [
                    {"cell_id": f"r{row + 1}c{column + 1}", "text": ""}
                    for column in range(table.data.num_cols)
                ]
                for row in range(table.data.num_rows)
            ]
            cell_bboxes: dict[str, list[float]] = {}
            for cell in table.data.table_cells:
                cell_id = f"r{cell.start_row_offset_idx + 1}c{cell.start_col_offset_idx + 1}"
                if (
                    cell.start_row_offset_idx < len(rows)
                    and cell.start_col_offset_idx < len(rows[cell.start_row_offset_idx])
                ):
                    rows[cell.start_row_offset_idx][cell.start_col_offset_idx] = {
                        "cell_id": cell_id,
                        "text": cell.text.strip(),
                    }
                if cell.bbox is not None and page is not None:
                    bbox = cell.bbox.to_bottom_left_origin(page.size.height)
                    cell_bboxes[cell_id] = [bbox.l, bbox.b, bbox.r, bbox.t]
            markdown = table.export_to_markdown(document).strip()
            if not markdown or not cell_bboxes:
                continue
            preview = _docling_preview_region(
                document, table, table_page, "table"
            )
            if preview is not None and isinstance(page_number, int):
                preview["page_number"] = page_number
            yield CandidateDraft(
                source_region_ids=(request.region_id,),
                channel_id="table",
                output_contract_version="eir-draft-v1",
                candidate_payload_ref=context.artifact_broker.put_text(markdown),
                content_kind_hint="table",
                element_kind_hint="table",
                table_grid={"table_id": f"table-{table_index}", "rows": rows},
                cell_bboxes=cell_bboxes,
                preview_region=preview,
                quality_flag_refs=(
                    () if preview else ("pdf_preview_region_missing",)
                ),
            )

        for picture in document.pictures:
            picture_page = next((value.page_no for value in picture.prov), None)
            if not isinstance(picture_page, int) or (
                isinstance(artifact_page_number, int)
                and picture_page != artifact_page_number
            ):
                continue
            caption = picture.caption_text(document).strip()
            if not caption:
                continue
            preview = _docling_preview_region(
                document, picture, picture_page, "figure"
            )
            if preview is not None and isinstance(page_number, int):
                preview["page_number"] = page_number
            yield CandidateDraft(
                source_region_ids=(request.region_id,),
                channel_id="generic_text",
                output_contract_version="eir-draft-v1",
                candidate_payload_ref=context.artifact_broker.put_text(caption),
                content_kind_hint="text",
                element_kind_hint="image",
                preview_region=preview,
                quality_flag_refs=(
                    () if preview else ("pdf_preview_region_missing",)
                ),
            )
