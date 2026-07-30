from __future__ import annotations

import base64
from datetime import datetime, timedelta, timezone
from io import BytesIO
import json
import asyncio
import os
from pathlib import Path
import sys
from types import SimpleNamespace
import pytest

from cryptography.hazmat.primitives import serialization
from cryptography.hazmat.primitives.asymmetric.ed25519 import Ed25519PrivateKey
from atlas_processing_sdk.package import build_package
from atlas_processing_sdk.scaffold import init_project

from fastapi.testclient import TestClient
from pypdf import PdfReader, PdfWriter

from atlas_plugin_runner import app as runner_app
from atlas_plugin_runner.app import create_app, _safe_child_diagnostic
from atlas_plugin_runner import host
from atlas_plugin_runner.builtin_plugins import (
    CsvPlugin,
    DoclingLayoutPlugin,
    DocxPlugin,
    InlineTextPlugin,
    PypdfPlugin,
    PptxPlugin,
    RapidOcrPlugin,
    XlsxPlugin,
    _docling_preview_region,
)


def tiny_png() -> bytes:
    from PIL import Image

    output = BytesIO()
    Image.new("RGB", (24, 12), "white").save(output, format="PNG")
    return output.getvalue()


def searchable_pdf(text: str) -> bytes:
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /Resources << /Font << /F1 4 0 R >> >> /MediaBox [0 0 612 792] /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    body = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, item in enumerate(objects, 1):
        offsets.append(len(body)); body.extend(f"{index} 0 obj\n".encode()); body.extend(item); body.extend(b"\nendobj\n")
    xref = len(body); body.extend(f"xref\n0 {len(objects)+1}\n".encode()); body.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]: body.extend(f"{offset:010d} 00000 n \n".encode())
    body.extend(f"trailer << /Size {len(objects)+1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode())
    return bytes(body)


def multipage_pdf(*page_texts: str | None) -> bytes:
    writer = PdfWriter()
    for text in page_texts:
        if text is None:
            writer.add_blank_page(width=612, height=792)
        else:
            writer.add_page(PdfReader(BytesIO(searchable_pdf(text))).pages[0])
    output = BytesIO()
    writer.write(output)
    return output.getvalue()


def empty_password_encrypted_pdf(text: str) -> bytes:
    writer = PdfWriter()
    writer.add_page(PdfReader(BytesIO(searchable_pdf(text))).pages[0])
    writer.encrypt(user_password="", owner_password="owner-password")
    output = BytesIO()
    writer.write(output)
    return output.getvalue()


def table_pdf() -> bytes:
    commands = [
        "2 w", "80 680 m 520 680 l S", "80 580 m 520 580 l S",
        "80 480 m 520 480 l S", "80 380 m 520 380 l S",
        "80 380 m 80 680 l S", "300 380 m 300 680 l S", "520 380 m 520 680 l S",
    ]
    for x, y, text in [
        (130, 625, "Part"), (350, 625, "Value"),
        (130, 525, "R1"), (350, 525, "10k"),
        (130, 425, "C1"), (350, 425, "100nF"),
    ]:
        commands.append(f"BT /F1 18 Tf {x} {y} Td ({text}) Tj ET")
    stream = "\n".join(commands).encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /Resources << /Font << /F1 4 0 R >> >> /MediaBox [0 0 612 792] /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    body = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, item in enumerate(objects, 1):
        offsets.append(len(body)); body.extend(f"{index} 0 obj\n".encode()); body.extend(item); body.extend(b"\nendobj\n")
    xref = len(body); body.extend(f"xref\n0 {len(objects)+1}\n".encode()); body.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]: body.extend(f"{offset:010d} 00000 n \n".encode())
    body.extend(f"trailer << /Size {len(objects)+1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode())
    return bytes(body)


def parser_request():
    return {
        "run_id": "run-1", "invocation_id": "inv-1", "document_id": "doc-1",
        "document_version_id": "dver-1", "artifact_ref": "artifact:1",
        "media_type": "application/pdf", "profile_id": "default-pdf",
        "profile_revision": 1, "policy_snapshot_ref": "policy:1",
        "deadline_at": (datetime.now(timezone.utc) + timedelta(minutes=1)).isoformat(),
        "batch_id": "batch-1", "unit_start": 1, "unit_end": 1,
        "resume_cursor": None,
        "plugin_config": {},
    }


def region_request():
    request = parser_request()
    for field in ("batch_id", "unit_start", "unit_end", "resume_cursor"):
        request.pop(field)
    return request


def invoke(client: TestClient, payload: dict):
    metadata = {
        key: value
        for key, value in payload.items()
        if key not in {"artifact", "package"}
    }
    files = [
        ("metadata", ("metadata.json", json.dumps(metadata), "application/json")),
        (
            "artifact",
            ("artifact.bin", payload["artifact"], "application/octet-stream"),
        ),
    ]
    package = payload.get("package")
    if package is not None:
        files.append(
            (
                "package",
                ("plugin.atlas-plugin", package, "application/octet-stream"),
            )
        )
    return client.post("/internal/v1/invocations", files=files)


def test_seccomp_is_loaded_before_external_module_import(monkeypatch, tmp_path):
    calls: list[str] = []
    monkeypatch.chdir(tmp_path)

    class Plugin:
        def __init__(self):
            calls.append("construct")

        async def parse(self, _request, _context):
            if False:
                yield None

    monkeypatch.setattr(host, "_apply_filesystem_isolation", lambda _payload: calls.append("landlock"))
    monkeypatch.setattr(host, "_deny_network_syscalls", lambda: calls.append("seccomp"))
    monkeypatch.setattr(
        host.importlib,
        "import_module",
        lambda _name: calls.append("import") or SimpleNamespace(Plugin=Plugin),
    )
    artifact_path = tmp_path / "artifact.bin"
    artifact_path.write_bytes(b"source")
    result = asyncio.run(host.execute({
        "entrypoint": "external_plugin:Plugin",
        "kind": "base_parser",
        "request": parser_request(),
        "artifact_path": str(artifact_path),
    }))
    assert result == {"drafts": [], "assets": {}}
    assert calls == ["landlock", "seccomp", "import", "construct"]


def unsigned_package(tmp_path: Path, monkeypatch, source: str) -> bytes:
    project = tmp_path / "plugin"
    init_project(project, "com.example.hostile", "base_parser")
    (project / "src" / "plugin.py").write_text(source)
    package = tmp_path / "hostile.atlas-plugin"
    build_package(project, package)
    monkeypatch.setenv("ATLAS_ALLOW_UNSIGNED_PLUGINS", "true")
    return package.read_bytes()


def external_payload(package: bytes, *, timeout_seconds: int = 3) -> dict:
    return {
        "invocation_id": "inv-hostile", "runtime_profile": "atlas-python-v1",
        "kind": "base_parser", "entrypoint": "plugin:Plugin",
        "request": parser_request(), "artifact": b"source",
        "package": package,
        "timeout_seconds": timeout_seconds,
    }


def test_network_attempt_is_denied_with_typed_safe_error(tmp_path, monkeypatch):
    package = unsigned_package(tmp_path, monkeypatch, """
import socket
socket.create_connection(("127.0.0.1", 9))
class Plugin:
    async def parse(self, request, context):
        if False: yield None
""")
    result = invoke(TestClient(create_app()), external_payload(package)).json()
    assert result == {"ok": False, "error": {"code": "plugin_network_denied", "type": "PermissionError"}}


def test_plugin_exception_is_typed_without_message_or_trace(tmp_path, monkeypatch):
    package = unsigned_package(tmp_path, monkeypatch, """
class Plugin:
    async def parse(self, request, context):
        raise RuntimeError("secret customer content")
        yield
""")
    result = invoke(TestClient(create_app()), external_payload(package)).json()
    assert result == {"ok": False, "error": {"code": "plugin_execution_failed", "type": "RuntimeError"}}


def test_plugin_keyboard_interrupt_is_an_execution_carrier_interruption(tmp_path, monkeypatch):
    package = unsigned_package(tmp_path, monkeypatch, """
class Plugin:
    async def parse(self, request, context):
        raise KeyboardInterrupt("do not expose this message")
        yield
""")
    result = invoke(TestClient(create_app()), external_payload(package)).json()
    assert result == {
        "ok": False,
        "error": {"code": "plugin_interrupted", "type": "KeyboardInterrupt"},
    }


def test_timeout_terminates_plugin_process(tmp_path, monkeypatch):
    package = unsigned_package(tmp_path, monkeypatch, """
import asyncio
class Plugin:
    async def parse(self, request, context):
        await asyncio.sleep(30)
        if False: yield None
""")
    result = invoke(
        TestClient(create_app()), external_payload(package, timeout_seconds=1)
    ).json()
    assert result == {"ok": False, "error": {"code": "plugin_timeout"}}


def test_raw_path_output_is_rejected(tmp_path, monkeypatch):
    package = unsigned_package(tmp_path, monkeypatch, """
from atlas_processing_sdk import SourceRegionDraft
class Plugin:
    async def parse(self, request, context):
        yield SourceRegionDraft(
            source_region_identity="page:1", region_kind="page",
            content_kind_hint="text", locator_draft={"page_number": 1},
            native_artifact_ref="/etc/passwd",
        )
""")
    result = invoke(TestClient(create_app()), external_payload(package)).json()
    assert result == {
        "ok": False,
        "error": {"code": "plugin_execution_failed", "type": "ValueError"},
    }


def test_minimal_environment_does_not_expose_runner_credentials(tmp_path, monkeypatch):
    monkeypatch.setenv("DATABASE_URL", "do-not-leak")
    package = unsigned_package(tmp_path, monkeypatch, """
import os
from atlas_processing_sdk import SourceRegionDraft
class Plugin:
    async def parse(self, request, context):
        leaked = os.environ.get("DATABASE_URL")
        yield SourceRegionDraft(
            source_region_identity="page:1", region_kind="page",
            content_kind_hint="text", locator_draft={"page_number": 1},
            normalized_text_ref="credential-present" if leaked else None,
        )
""")
    result = invoke(TestClient(create_app()), external_payload(package)).json()
    assert result["ok"] is True, result
    assert result["drafts"][0]["normalized_text_ref"] is None


def test_external_plugin_object_state_is_fresh_for_every_logical_invocation(
    tmp_path, monkeypatch,
):
    package = unsigned_package(tmp_path, monkeypatch, """
from atlas_processing_sdk import SourceRegionDraft
counter = 0
class Plugin:
    async def parse(self, request, context):
        global counter
        counter += 1
        yield SourceRegionDraft(
            source_region_identity="page:1", region_kind="page",
            content_kind_hint="text",
            locator_draft={"page_number": counter},
        )
""")
    client = TestClient(create_app())
    first = invoke(
        client,
        {**external_payload(package), "invocation_id": "inv-state-first"},
    ).json()
    second = invoke(
        client,
        {**external_payload(package), "invocation_id": "inv-state-second"},
    ).json()

    assert first["ok"] is True, first
    assert second["ok"] is True, second
    assert first["drafts"][0]["locator_draft"]["page_number"] == 1
    assert second["drafts"][0]["locator_draft"]["page_number"] == 1


@pytest.mark.skipif(not sys.platform.startswith("linux"), reason="Landlock is a Linux kernel boundary")
def test_landlock_blocks_etc_passwd_before_plugin_import(tmp_path, monkeypatch):
    package = unsigned_package(tmp_path, monkeypatch, """
open("/etc/passwd", encoding="utf-8").read()
class Plugin:
    async def parse(self, request, context):
        if False: yield None
""")
    result = invoke(TestClient(create_app()), external_payload(package)).json()
    assert result == {"ok": False, "error": {"code": "plugin_filesystem_denied", "type": "PermissionError"}}


@pytest.mark.skipif(not sys.platform.startswith("linux"), reason="seccomp is a Linux kernel boundary")
def test_seccomp_allows_threads_but_denies_child_processes(tmp_path, monkeypatch):
    threaded = unsigned_package(tmp_path / "threaded", monkeypatch, """
import threading
from atlas_processing_sdk import SourceRegionDraft
class Plugin:
    async def parse(self, request, context):
        values = []
        worker = threading.Thread(target=lambda: values.append("ok"))
        worker.start(); worker.join()
        yield SourceRegionDraft(
            source_region_identity="page:1", region_kind="page",
            content_kind_hint="text", locator_draft={"page_number": 1},
            normalized_text_ref="opaque:thread-ok" if values == ["ok"] else None,
        )
""")
    allowed = invoke(TestClient(create_app()), external_payload(threaded)).json()
    assert allowed["ok"] is True, allowed

    forking = unsigned_package(tmp_path / "forking", monkeypatch, """
import os
os.fork()
class Plugin:
    async def parse(self, request, context):
        if False: yield None
""")
    denied = invoke(TestClient(create_app()), external_payload(forking)).json()
    assert denied == {
        "ok": False,
        "error": {"code": "plugin_filesystem_denied", "type": "PermissionError"},
    }


def test_real_docling_emits_generic_text_and_grounded_table(tmp_path, monkeypatch):
    model_root = Path(os.environ.get("ATLAS_DOCLING_ARTIFACTS_PATH", "/opt/docling-models"))
    if not model_root.exists():
        pytest.skip("Docling model cache is unavailable")
    project = Path(__file__).resolve().parents[2] / "plugin-sdk" / "examples" / "atlas-docling-pdf"
    package_path = tmp_path / "docling.atlas-plugin"
    build_package(project, package_path)
    monkeypatch.setenv("ATLAS_ALLOW_UNSIGNED_PLUGINS", "true")
    request = {
        **region_request(), "region_id": "region-page-1", "region_kind": "page",
        "element_kind_hint": "page", "content_kind_hint": "text",
        "normalized_text_ref": None, "structured_content_ref": None,
        "locator_draft": {"selector_kind": "page_region", "page_number": 1},
        "active_trait_hints": [],
    }
    payload = external_payload(package_path.read_bytes(), timeout_seconds=120)
    payload.update({
        "runtime_profile": "atlas-docling-cpu-v1", "kind": "region_processor",
        "request": request, "artifact": table_pdf(),
    })
    result = invoke(TestClient(create_app()), payload).json()
    assert result["ok"] is True, result
    channels = {draft["channel_id"] for draft in result["drafts"]}
    assert {"generic_text", "table"} <= channels
    table = next(draft for draft in result["drafts"] if draft["channel_id"] == "table")
    assert table["preview_region"]["region_kind"] == "table"
    assert table["preview_region"]["coordinate_system"] == "pdf_crop_box_relative_bottom_left"
    assert table["preview_region"]["rectangles"]
    assert table["table_grid"]["rows"] == [
        [{"cell_id": "r1c1", "text": "Part"}, {"cell_id": "r1c2", "text": "Value"}],
        [{"cell_id": "r2c1", "text": "R1"}, {"cell_id": "r2c2", "text": "10k"}],
        [{"cell_id": "r3c1", "text": "C1"}, {"cell_id": "r3c2", "text": "100nF"}],
    ]
    assert set(table["cell_bboxes"]) == {"r1c1", "r1c2", "r2c1", "r2c2", "r3c1", "r3c2"}
    paragraphs = [
        draft for draft in result["drafts"]
        if draft["channel_id"] == "generic_text"
        and draft["element_kind_hint"] == "paragraph"
    ]
    assert paragraphs and all(draft["preview_region"] for draft in paragraphs)


async def run_fake_host(
    tmp_path: Path,
    body: str,
    disconnected=None,
    host_payload: dict | None = None,
):
    script = tmp_path / "fake_host.py"
    script.write_text(body)
    async def connected(): return False
    return await runner_app._run_host(
        script, tmp_path, host_payload or {}, 3, disconnected or connected,
    )


def test_child_disconnect_is_cancelled(tmp_path):
    calls = 0
    async def disconnected():
        nonlocal calls
        calls += 1
        return calls > 1
    output, error = asyncio.run(run_fake_host(
        tmp_path, "import sys,time; sys.stdin.read(); time.sleep(30)", disconnected,
    ))
    assert output is None
    assert error == "plugin_cancelled"


def test_child_crash_is_typed(tmp_path):
    output, error = asyncio.run(run_fake_host(tmp_path, "import os,sys; sys.stdin.read(); os._exit(9)"))
    assert output is None
    assert error == "plugin_crashed"


def test_docling_child_receives_bounded_cpu_thread_environment(tmp_path):
    output, error = asyncio.run(run_fake_host(
        tmp_path,
        """
import json, os, sys
json.load(sys.stdin)
print(json.dumps({name: os.environ.get(name) for name in (
    'OMP_NUM_THREADS', 'OMP_THREAD_LIMIT', 'OMP_STACKSIZE',
    'OPENBLAS_NUM_THREADS', 'MKL_NUM_THREADS', 'KMP_STACKSIZE',
    'MALLOC_ARENA_MAX', 'DOCLING_NUM_THREADS', 'DOCLING_DEVICE',
)}))
""",
        host_payload={"runtime_profile": "atlas-docling-cpu-v1"},
    ))
    assert error is None
    assert json.loads(output) == {
        "OMP_NUM_THREADS": "4",
        "OMP_THREAD_LIMIT": "8",
        "OMP_STACKSIZE": "1M",
        "OPENBLAS_NUM_THREADS": "1",
        "MKL_NUM_THREADS": "4",
        "KMP_STACKSIZE": "1m",
        "MALLOC_ARENA_MAX": "2",
        "DOCLING_NUM_THREADS": "4",
        "DOCLING_DEVICE": "cpu",
    }


def test_child_stderr_is_reduced_to_safe_diagnostic_categories():
    assert _safe_child_diagnostic(
        b"libgomp: Thread creation failed: Resource temporarily unavailable"
    ) == "thread_creation_failed"
    assert _safe_child_diagnostic(b"Stage preprocess: std::bad_alloc") == "memory_limit"
    assert _safe_child_diagnostic(b"confidential document content") == "process_exit"


def test_child_crash_log_does_not_include_stderr_content(tmp_path, caplog):
    with caplog.at_level("WARNING", logger="atlas_plugin_runner"):
        output, error = asyncio.run(run_fake_host(
            tmp_path,
            """
import os, sys
sys.stdin.read()
sys.stderr.write('confidential document content')
sys.stderr.flush()
os._exit(7)
""",
        ))
    assert output is None
    assert error == "plugin_crashed"
    assert "category=process_exit" in caplog.text
    assert "confidential document content" not in caplog.text


def test_non_json_output_is_protocol_error(tmp_path):
    output, error = asyncio.run(run_fake_host(tmp_path, "import sys; sys.stdin.read(); print('not-json')"))
    assert error is None
    assert output == "not-json\n"
    try:
        json.loads(output)
    except json.JSONDecodeError:
        pass
    else:
        raise AssertionError("fake host output unexpectedly parsed")


def test_api_maps_non_json_output_to_protocol_error(monkeypatch):
    async def fake_run(*_args, **_kwargs): return "not-json", None
    monkeypatch.setattr(runner_app, "_run_host", fake_run)
    result = invoke(TestClient(create_app()), {
        "invocation_id": "inv-protocol", "runtime_profile": "atlas-python-v1",
        "kind": "base_parser", "entrypoint": "atlas_plugin_runner.builtin_plugins:PypdfPlugin",
        "request": parser_request(), "artifact": b"source",
    }).json()
    assert result == {"ok": False, "error": {"code": "runner_protocol_error"}}


def test_oversized_stdout_and_stderr_are_terminated(tmp_path, monkeypatch):
    monkeypatch.setattr(runner_app, "STDOUT_LIMIT", 32)
    output, error = asyncio.run(run_fake_host(
        tmp_path, "import sys; sys.stdin.read(); sys.stdout.write('x'*10000); sys.stdout.flush()",
    ))
    assert output is None
    assert error == "plugin_output_limit_exceeded"
    monkeypatch.setattr(runner_app, "STDOUT_LIMIT", 64 * 1024 * 1024)
    monkeypatch.setattr(runner_app, "STDERR_LIMIT", 32)
    output, error = asyncio.run(run_fake_host(
        tmp_path, "import sys; sys.stdin.read(); sys.stderr.write('x'*10000); sys.stderr.flush()",
    ))
    assert output is None
    assert error == "plugin_output_limit_exceeded"


def test_landlock_unavailable_fails_closed(monkeypatch):
    monkeypatch.setattr(host.sys, "platform", "linux")
    monkeypatch.setattr(host, "_landlock_syscalls", lambda: (1, 2, 3))
    class Libc:
        class Syscall:
            restype = None
            def __call__(self, *_args): return -1
        syscall = Syscall()
    monkeypatch.setattr(host.ctypes, "CDLL", lambda *_args, **_kwargs: Libc())
    try:
        host._apply_filesystem_isolation({})
    except RuntimeError as exc:
        assert "Landlock" in str(exc)
    else:
        raise AssertionError("missing Landlock did not fail closed")


def test_runner_derives_trusted_builtin_only_after_request_validation(monkeypatch):
    captured = {}

    async def fake_run(_host_path, _workspace, payload, *_args):
        captured.update(payload)
        return json.dumps({"ok": True, "drafts": [], "assets": {}}), None

    monkeypatch.setattr(runner_app, "_run_host", fake_run)
    client = TestClient(create_app())
    request = {
        "invocation_id": "inv-runtime-roots",
        "runtime_profile": "atlas-python-v1",
        "kind": "base_parser",
        "entrypoint": "atlas_plugin_runner.builtin_plugins:PypdfPlugin",
        "request": parser_request(),
        "artifact": b"source",
    }
    response = invoke(client, request)
    assert response.status_code == 200
    assert captured["trusted_builtin"] is True
    assert {
        field: captured["request"][field]
        for field in ("batch_id", "unit_start", "unit_end", "resume_cursor")
    } == {
        "batch_id": "batch-1", "unit_start": 1, "unit_end": 1,
        "resume_cursor": None,
    }
    assert captured["artifact_path"].endswith("/workspace/artifact.bin")
    assert not Path(captured["artifact_path"]).exists()
    assert "artifact_base64" not in captured
    assert "package_base64" not in captured

    rejected = invoke(client, {**request, "trusted_builtin": True})
    assert rejected.status_code == 422


def test_runner_http_contract_accepts_the_trusted_docling_processor(monkeypatch):
    captured = {}

    async def fake_run(_host_path, _workspace, payload, *_args):
        captured.update(payload)
        return json.dumps({"ok": True, "drafts": [], "assets": {}}), None

    monkeypatch.setattr(runner_app, "_run_host", fake_run)
    request = {
        **region_request(),
        "region_id": "page:1",
        "region_kind": "page",
        "content_kind_hint": "unknown",
        "element_kind_hint": "page",
        "normalized_text_ref": None,
        "structured_content_ref": None,
        "locator_draft": {"selector_kind": "page_region", "page_number": 1},
        "active_trait_hints": [],
    }
    result = invoke(TestClient(create_app()), {
        "invocation_id": "inv-docling-http",
        "runtime_profile": "atlas-docling-cpu-v1",
        "kind": "region_processor",
        "entrypoint": "atlas_plugin_runner.builtin_plugins:DoclingLayoutPlugin",
        "request": request,
        "artifact": searchable_pdf("Docling contract"),
    }).json()

    assert result["ok"] is True
    assert captured["trusted_builtin"] is True
    assert captured["runtime_profile"] == "atlas-docling-cpu-v1"


def test_runner_rejects_legacy_json_artifact_envelope():
    response = TestClient(create_app()).post(
        "/internal/v1/invocations",
        json={
            "invocation_id": "inv-legacy-json",
            "runtime_profile": "atlas-python-v1",
            "kind": "base_parser",
            "entrypoint": "atlas_plugin_runner.builtin_plugins:PypdfPlugin",
            "request": parser_request(),
            "artifact_base64": base64.b64encode(b"source").decode(),
        },
    )
    assert response.status_code == 200
    assert response.json() == {
        "ok": False,
        "error": {"code": "invalid_artifact_envelope"},
    }


def test_runner_rejects_oversized_multipart_artifact_before_host(monkeypatch):
    host_called = False

    async def fake_run(*_args, **_kwargs):
        nonlocal host_called
        host_called = True
        return json.dumps({"ok": True, "drafts": [], "assets": {}}), None

    monkeypatch.setattr(runner_app, "ARTIFACT_LIMIT", 4)
    monkeypatch.setattr(runner_app, "_run_host", fake_run)
    response = invoke(
        TestClient(create_app()),
        {
            "invocation_id": "inv-too-large",
            "runtime_profile": "atlas-python-v1",
            "kind": "base_parser",
            "entrypoint": "atlas_plugin_runner.builtin_plugins:PypdfPlugin",
            "request": parser_request(),
            "artifact": b"12345",
        },
    )
    assert response.json() == {
        "ok": False,
        "error": {"code": "invalid_artifact_envelope"},
    }
    assert host_called is False


def test_builtin_pypdf_and_generic_processor_round_trip():
    client = TestClient(create_app())
    artifact = searchable_pdf("Atlas runner text")
    parsed = invoke(client, {
        "invocation_id": "inv-1", "runtime_profile": "atlas-python-v1",
        "kind": "base_parser", "entrypoint": "atlas_plugin_runner.builtin_plugins:PypdfPlugin",
        "request": parser_request(), "artifact": artifact,
    })
    assert parsed.status_code == 200
    payload = parsed.json()
    assert payload["ok"] is True, payload
    assert payload["drafts"][0]["element_kind_hint"] == "page"
    assert payload["drafts"][1]["element_kind_hint"] == "paragraph"
    ref = payload["drafts"][0]["normalized_text_ref"]
    region_payload = {**region_request(),
        "region_id": "region-1", "region_kind": "page", "element_kind_hint": "page",
        "content_kind_hint": "text", "normalized_text_ref": ref,
        "structured_content_ref": None, "locator_draft": {"page_number": 1},
        "active_trait_hints": [],
    }
    processed = invoke(client, {
        "invocation_id": "inv-2", "runtime_profile": "atlas-python-v1",
        "kind": "region_processor", "entrypoint": "atlas_plugin_runner.builtin_plugins:GenericTextPlugin",
        "request": region_payload, "artifact": artifact, "input_assets": payload["assets"],
    }).json()
    assert processed["ok"] is True
    assert processed["drafts"][0]["channel_id"] == "generic_text"


def test_builtin_pypdf_parses_empty_password_encrypted_source() -> None:
    parsed = invoke(TestClient(create_app()), {
        "invocation_id": "inv-empty-password",
        "runtime_profile": "atlas-python-v1",
        "kind": "base_parser",
        "entrypoint": "atlas_plugin_runner.builtin_plugins:PypdfPlugin",
        "request": parser_request(),
        "artifact": empty_password_encrypted_pdf("Permission protected text"),
    })

    assert parsed.status_code == 200
    payload = parsed.json()
    assert payload["ok"] is True, payload
    assert any(
        value == base64.b64encode(b"Permission protected text").decode()
        for value in payload["assets"].values()
    )


def test_pypdf_deterministically_emits_multipage_paragraph_regions():
    class Broker:
        def __init__(self): self.outputs = {}
        async def parsed_pdf_pages(self, _ref, unit_start, unit_end):
            assert (unit_start, unit_end) == (2, 3)
            return [(2, "text:2"), (3, "text:3")]
        async def read_text(self, ref):
            return {"text:2": "First paragraph.\n\nSecond paragraph.", "text:3": "Third paragraph."}[ref]
        def put_text(self, value):
            ref = f"opaque:{len(self.outputs) + 1}"
            self.outputs[ref] = value
            return ref

    broker = Broker()
    context = SimpleNamespace(artifact_broker=broker)

    async def collect():
        return [
            draft async for draft in PypdfPlugin().parse(
                SimpleNamespace(
                    artifact_ref="artifact:opaque", unit_start=2, unit_end=3
                ),
                context,
            )
        ]

    drafts = asyncio.run(collect())
    assert [(draft.region_kind, draft.locator_draft) for draft in drafts] == [
        ("page", {"selector_kind": "page_region", "page_number": 2}),
        ("paragraph", {"selector_kind": "normalized_text_span", "page_number": 2, "ordinal": 1}),
        ("paragraph", {"selector_kind": "normalized_text_span", "page_number": 2, "ordinal": 2}),
        ("page", {"selector_kind": "page_region", "page_number": 3}),
        ("paragraph", {"selector_kind": "normalized_text_span", "page_number": 3, "ordinal": 1}),
    ]
    assert [draft.normalized_text_ref for draft in drafts if draft.region_kind == "paragraph"] == [
        "opaque:1", "opaque:2", "opaque:3",
    ]


def test_office_image_processor_emits_local_ocr_channel(monkeypatch):
    import rapidocr

    class FakeRapidOCR:
        def __call__(self, _content):
            return SimpleNamespace(txts=("VISIBLE 47 OHM",))

    observed_params = []

    def fake_rapidocr(*, params):
        observed_params.append(params)
        return FakeRapidOCR()

    monkeypatch.setattr(rapidocr, "RapidOCR", fake_rapidocr)

    class Broker:
        def __init__(self):
            self.assets = {"asset:image": tiny_png()}
            self.outputs = {}

        async def read_bytes(self, ref):
            return self.assets[ref]

        def put_text(self, value):
            ref = f"text:{len(self.outputs) + 1}"
            self.outputs[ref] = value
            return ref

        def put_bytes(self, value):
            ref = f"bytes:{len(self.outputs) + 1}"
            self.outputs[ref] = value
            return ref

    request = SimpleNamespace(
        native_artifact_ref="asset:image",
        region_id="slide:1:image:2",
        locator_draft={"selector_kind": "powerpoint_image"},
    )
    broker = Broker()

    async def collect(plugin):
        return [
            draft async for draft in plugin.process(
                request, SimpleNamespace(artifact_broker=broker)
            )
        ]

    ocr = asyncio.run(collect(RapidOcrPlugin()))
    assert observed_params == [{
        "EngineConfig.onnxruntime.intra_op_num_threads": 1,
        "EngineConfig.onnxruntime.inter_op_num_threads": 1,
    }]
    assert ocr[0].channel_id == "generic_text"
    assert ocr[0].element_kind_hint == "ocr_text"
    assert broker.outputs[ocr[0].candidate_payload_ref] == "VISIBLE 47 OHM"


def test_docling_preview_region_uses_declared_bottom_left_page_geometry():
    transformed = SimpleNamespace(l=10, b=20, r=110, t=220)
    bbox = SimpleNamespace(
        to_bottom_left_origin=lambda page_height: transformed
        if page_height == 792 else None
    )
    item = SimpleNamespace(
        self_ref="#/texts/1",
        prov=[SimpleNamespace(page_no=1, bbox=bbox)],
    )
    document = SimpleNamespace(
        pages={1: SimpleNamespace(size=SimpleNamespace(width=612, height=792))}
    )
    for region_kind in ("paragraph", "table", "figure", "image"):
        assert _docling_preview_region(document, item, 1, region_kind) == {
            "page_number": 1,
            "region_kind": region_kind,
            "source_element_id": "#/texts/1",
            "coordinate_system": "pdf_crop_box_relative_bottom_left",
            "rectangles": [[10, 20, 110, 220]],
            "page_width": 612,
            "page_height": 792,
            "geometry_version": "docling-page-region-v1",
        }


def test_docling_caption_text_is_anchored_to_the_whole_figure(monkeypatch):
    from docling import document_converter
    from docling.datamodel.base_models import ConversionStatus

    transformed = SimpleNamespace(l=40, b=50, r=400, t=500)
    picture = SimpleNamespace(
        self_ref="#/pictures/1",
        prov=[SimpleNamespace(
            page_no=1,
            bbox=SimpleNamespace(
                to_bottom_left_origin=lambda _height: transformed
            ),
        )],
        caption_text=lambda _document: "Power stage efficiency chart",
    )
    document = SimpleNamespace(
        iterate_items=lambda: iter(()),
        tables=[],
        pictures=[picture],
        pages={1: SimpleNamespace(size=SimpleNamespace(width=612, height=792))},
    )

    class Converter:
        def __init__(self, **_kwargs):
            pass

        def convert(self, *_args, **_kwargs):
            return SimpleNamespace(
                status=ConversionStatus.SUCCESS,
                document=document,
            )

    monkeypatch.setattr(document_converter, "DocumentConverter", Converter)

    class Broker:
        def __init__(self):
            self.assets = {}

        async def read_bytes(self, _ref):
            return searchable_pdf("caption")

        def put_text(self, value):
            self.assets["runner-text:1"] = value
            return "runner-text:1"

    broker = Broker()

    async def collect():
        return [draft async for draft in DoclingLayoutPlugin().process(
            SimpleNamespace(
                artifact_ref="artifact:figure",
                locator_draft={"page_number": 1},
                region_id="page:1",
            ),
            SimpleNamespace(artifact_broker=broker),
        )]

    drafts = asyncio.run(collect())
    assert len(drafts) == 1
    assert broker.assets == {"runner-text:1": "Power stage efficiency chart"}
    assert drafts[0].element_kind_hint == "image"
    assert drafts[0].preview_region == {
        "page_number": 1,
        "region_kind": "figure",
        "source_element_id": "#/pictures/1",
        "coordinate_system": "pdf_crop_box_relative_bottom_left",
        "rectangles": [[40, 50, 400, 500]],
        "page_width": 612,
        "page_height": 792,
        "geometry_version": "docling-page-region-v1",
    }


def test_docling_partial_conversion_fails_closed(monkeypatch):
    from docling import document_converter
    from docling.datamodel.base_models import ConversionStatus

    class Converter:
        def __init__(self, **_kwargs):
            pass

        def convert(self, *_args, **_kwargs):
            return SimpleNamespace(
                status=ConversionStatus.PARTIAL_SUCCESS,
                document=SimpleNamespace(),
            )

    monkeypatch.setattr(document_converter, "DocumentConverter", Converter)

    class Broker:
        async def read_bytes(self, _ref):
            return searchable_pdf("partial")

    async def collect():
        return [draft async for draft in DoclingLayoutPlugin().process(
            SimpleNamespace(
                artifact_ref="artifact:partial",
                locator_draft={"page_number": 1},
                region_id="page:1",
            ),
            SimpleNamespace(artifact_broker=Broker()),
        )]

    with pytest.raises(RuntimeError, match="docling_conversion_incomplete"):
        asyncio.run(collect())


def test_builtin_pypdf_emits_only_selected_range_and_empty_page_is_deterministic():
    client = TestClient(create_app())
    artifact = multipage_pdf("Before range", None, "Selected page", "After range")
    request = {**parser_request(), "batch_id": "batch-2-3", "unit_start": 2, "unit_end": 3}
    invocation = {
        "invocation_id": "inv-range", "runtime_profile": "atlas-python-v1",
        "kind": "base_parser",
        "entrypoint": "atlas_plugin_runner.builtin_plugins:PypdfPlugin",
        "request": request, "artifact": artifact,
    }

    selected = invoke(client, invocation).json()
    replayed = invoke(client, invocation).json()
    assert selected == replayed
    assert selected["ok"] is True, selected
    assert [draft["source_region_identity"] for draft in selected["drafts"]] == [
        "page:2", "page:3", "page:3:paragraph:1",
    ]
    assert [draft["normalized_text_ref"] for draft in selected["drafts"]] == [
        None, "runner-text:1", "runner-text:2",
    ]
    decoded_assets = "\n".join(
        base64.b64decode(value).decode("utf-8")
        for value in selected["assets"].values()
    )
    assert "Selected page" in decoded_assets
    assert "Before range" not in decoded_assets
    assert "After range" not in decoded_assets

    empty_request = {
        **request, "batch_id": "batch-empty", "unit_start": 2, "unit_end": 2
    }
    empty_invocation = {**invocation, "request": empty_request}
    empty = invoke(client, empty_invocation).json()
    assert empty == invoke(client, empty_invocation).json()
    assert empty == {
        "ok": True,
        "drafts": [{
            "source_region_identity": "page:2",
            "region_kind": "page",
            "content_kind_hint": "unknown",
            "element_kind_hint": "page",
            "locator_draft": {"selector_kind": "page_region", "page_number": 2},
            "parent_region_identity": None,
            "normalized_text_ref": None,
            "structured_content_ref": None,
            "native_artifact_ref": None,
            "quality_flag_refs": [],
        }],
        "assets": {},
    }


def test_inline_text_parser_accepts_batch_contract_without_page_state():
    result = invoke(TestClient(create_app()), {
        "invocation_id": "inv-inline", "runtime_profile": "atlas-python-v1",
        "kind": "base_parser",
        "entrypoint": "atlas_plugin_runner.builtin_plugins:InlineTextPlugin",
        "request": {
            **parser_request(), "media_type": "text/plain", "batch_id": "batch-text",
            "unit_start": 7, "unit_end": 9, "resume_cursor": "cursor:text-7",
        },
        "artifact": b"First paragraph.\n\nSecond paragraph.",
    }).json()

    assert result["ok"] is True, result
    assert [draft["source_region_identity"] for draft in result["drafts"]] == [
        "paragraph:1", "paragraph:2",
    ]
    assert all("page_number" not in draft["locator_draft"] for draft in result["drafts"])


class _NativeParserBroker:
    def __init__(self, source: bytes):
        self.source = source
        self.assets: dict[str, bytes | str] = {}

    async def read_bytes(self, _ref):
        return self.source

    def put_text(self, value):
        ref = f"runner-text:{len(self.assets) + 1}"
        self.assets[ref] = value
        return ref

    def put_bytes(self, value):
        ref = f"runner-bytes:{len(self.assets) + 1}"
        self.assets[ref] = value
        return ref


def _collect_native(parser, source: bytes):
    broker = _NativeParserBroker(source)

    async def collect():
        return [
            item async for item in parser.parse(
                SimpleNamespace(artifact_ref="artifact:native"),
                SimpleNamespace(artifact_broker=broker),
            )
        ]

    return asyncio.run(collect()), broker


def test_txt_and_csv_parsers_preserve_line_and_cell_identities():
    text_drafts, text_broker = _collect_native(
        InlineTextPlugin(), b"First line\ncontinued\n\nSecond paragraph"
    )
    assert [draft.locator_draft for draft in text_drafts] == [
        {"selector_kind": "text_line_span", "ordinal": 1, "line_start": 1, "line_end": 2},
        {"selector_kind": "text_line_span", "ordinal": 2, "line_start": 4, "line_end": 4},
    ]
    assert set(text_broker.assets.values()) == {"First line\ncontinued", "Second paragraph"}

    csv_drafts, csv_broker = _collect_native(
        CsvPlugin(), 'name,value\n"R,1",10k\nC1,100nF'.encode()
    )
    assert [draft.locator_draft["row_number"] for draft in csv_drafts] == [1, 2, 3]
    assert csv_drafts[1].content_kind_hint == "table"
    structured = json.loads(csv_broker.assets[csv_drafts[1].structured_content_ref])
    assert structured["cells"] == [
        {"column_index": 1, "value": "R,1"},
        {"column_index": 2, "value": "10k"},
    ]


def test_docx_parser_preserves_each_image_occurrence_with_text_anchors():
    from docx import Document
    from PIL import Image

    image = BytesIO()
    Image.new("RGB", (3, 2), "red").save(image, format="PNG")
    image.seek(0)
    document = Document()
    document.add_heading("Power Design", level=1)
    document.add_paragraph("Before first image")
    document.add_picture(BytesIO(image.getvalue()))
    document.add_paragraph("Between duplicate images")
    document.add_picture(BytesIO(image.getvalue()))
    document.add_paragraph("After second image")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Part"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "R1"
    table.cell(1, 1).text = "10k"
    output = BytesIO()
    document.save(output)

    drafts, broker = _collect_native(DocxPlugin(), output.getvalue())
    assert any(draft.element_kind_hint == "heading" for draft in drafts)
    word_table = next(draft for draft in drafts if draft.locator_draft["selector_kind"] == "word_table")
    assert "R1\t10k" in broker.assets[word_table.normalized_text_ref]
    embedded = [
        draft for draft in drafts
        if draft.locator_draft["selector_kind"] == "word_image"
    ]
    assert len(embedded) == 2
    assert embedded[0].locator_draft["relationship_id"] == embedded[1].locator_draft["relationship_id"]
    assert [draft.locator_draft["image_index"] for draft in embedded] == [1, 2]
    assert [draft.locator_draft["alignment_anchors"] for draft in embedded] == [
        ["Before first image", "Between duplicate images"],
        ["Between duplicate images", "After second image"],
    ]
    assert all(
        bytes(broker.assets[draft.native_artifact_ref]).startswith(b"\x89PNG")
        for draft in embedded
    )
    assert {draft.locator_draft["part_name"] for draft in embedded} == {
        "word/media/image1.png"
    }


def test_xlsx_parser_preserves_sheet_rows_table_and_anchored_image():
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as SpreadsheetImage
    from openpyxl.worksheet.table import Table
    from PIL import Image

    image = BytesIO()
    Image.new("RGB", (3, 2), "blue").save(image, format="PNG")
    image.seek(0)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "BOM"
    sheet.append(["Part", "Value"])
    sheet.append(["R1", "10k"])
    sheet.append([None, None])
    sheet.append(["R2", "22k"])
    sheet.add_table(Table(displayName="BomTable", ref="A1:B4"))
    sheet.add_image(SpreadsheetImage(image), "D3")
    output = BytesIO()
    workbook.save(output)

    drafts, broker = _collect_native(XlsxPlugin(), output.getvalue())
    row = next(draft for draft in drafts if draft.locator_draft.get("row_index") == 2)
    assert broker.assets[row.normalized_text_ref] == "R1\t10k"
    assert any(draft.locator_draft.get("table_name") == "BomTable" for draft in drafts)
    embedded = next(draft for draft in drafts if draft.locator_draft["selector_kind"] == "excel_image")
    assert embedded.locator_draft["anchor_row"] == 3
    assert embedded.locator_draft["alignment_anchors"] == ["R1\t10k", "R2\t22k"]
    assert bytes(broker.assets[embedded.native_artifact_ref]).startswith(b"\x89PNG")


def test_pptx_parser_preserves_slide_shape_chart_table_and_image(monkeypatch):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import pptx

    text_shape = SimpleNamespace(
        shape_type=MSO_SHAPE_TYPE.TEXT_BOX, has_text_frame=True,
        has_table=False, has_chart=False, text="Architecture", name="Title",
        left=1, top=2, width=3, height=4,
    )
    table_shape = SimpleNamespace(
        shape_type=MSO_SHAPE_TYPE.TABLE, has_text_frame=False,
        has_table=True, has_chart=False, name="Parts", left=5, top=6, width=7, height=8,
        table=SimpleNamespace(rows=[SimpleNamespace(cells=[SimpleNamespace(text="R1"), SimpleNamespace(text="10k")])]),
    )
    series = SimpleNamespace(name="Efficiency", values=[90, 95])
    chart_shape = SimpleNamespace(
        shape_type=MSO_SHAPE_TYPE.CHART, has_text_frame=False,
        has_table=False, has_chart=True, name="Chart", left=9, top=10, width=11, height=12,
        chart=SimpleNamespace(has_title=False, series=[series]),
    )
    image_shape = SimpleNamespace(
        shape_type=MSO_SHAPE_TYPE.PICTURE, has_text_frame=False,
        has_table=False, has_chart=False, name="Screenshot", left=13, top=14,
        width=15, height=16, image=SimpleNamespace(blob=b"image-bytes"),
    )
    monkeypatch.setattr(
        pptx,
        "Presentation",
            lambda _stream: SimpleNamespace(
                slides=[SimpleNamespace(shapes=[text_shape, table_shape, chart_shape, image_shape])],
                slide_width=100,
                slide_height=80,
            ),
    )
    drafts, broker = _collect_native(PptxPlugin(), b"pptx-placeholder")
    selector_kinds = {draft.locator_draft["selector_kind"] for draft in drafts}
    assert {
        "powerpoint_slide", "powerpoint_shape", "powerpoint_table",
        "powerpoint_chart", "powerpoint_image",
    } <= selector_kinds
    image_draft = next(draft for draft in drafts if draft.locator_draft["selector_kind"] == "powerpoint_image")
    assert broker.assets[image_draft.native_artifact_ref] == b"image-bytes"


def test_runner_returns_safe_error_without_trace_or_content():
    result = invoke(TestClient(create_app()), {
        "invocation_id": "inv-bad", "runtime_profile": "atlas-python-v1",
        "kind": "base_parser", "entrypoint": "missing.module:Plugin",
        "request": parser_request(), "artifact": b"secret-content",
    }).json()
    assert result == {"ok": False, "error": {"code": "builtin_contract_mismatch"}}


def test_external_invocation_must_match_trusted_package_identity(tmp_path, monkeypatch):
    project = tmp_path / "plugin"
    init_project(project, "com.example.identity", "base_parser")
    private = Ed25519PrivateKey.generate()
    private_path = tmp_path / "private.pem"
    private_path.write_bytes(private.private_bytes(
        serialization.Encoding.PEM, serialization.PrivateFormat.PKCS8,
        serialization.NoEncryption(),
    ))
    public = private.public_key().public_bytes(
        serialization.Encoding.PEM, serialization.PublicFormat.SubjectPublicKeyInfo,
    ).decode()
    package = tmp_path / "identity.atlas-plugin"
    build_package(project, package, signing_key=private_path, signing_key_id="team-key")
    monkeypatch.setenv("ATLAS_PLUGIN_TRUSTED_KEYS_JSON", json.dumps({"team-key": public}))
    result = invoke(TestClient(create_app()), {
        "invocation_id": "inv-mismatch", "runtime_profile": "atlas-docling-cpu-v1",
        "kind": "base_parser", "entrypoint": "plugin:Plugin",
        "request": parser_request(), "artifact": b"source",
        "package": package.read_bytes(),
    }).json()
    assert result == {"ok": False, "error": {"code": "plugin_package_untrusted"}}


def test_external_invocation_enforces_packaged_output_schema(tmp_path, monkeypatch):
    project = tmp_path / "plugin"
    init_project(project, "com.example.schema", "base_parser")
    (project / "schemas" / "output.schema.json").write_text(json.dumps({
        "$schema": "https://json-schema.org/draft/2020-12/schema",
        "type": "object",
        "required": ["schema_required_marker"],
    }))
    (project / "src" / "plugin.py").write_text('''
from atlas_processing_sdk import SourceRegionDraft

class Plugin:
    async def parse(self, request, context):
        yield SourceRegionDraft(
            source_region_identity="page:1",
            region_kind="page",
            content_kind_hint="text",
            element_kind_hint="page",
            locator_draft={"selector_kind": "page_region", "page_number": 1},
        )
''')
    package = tmp_path / "schema.atlas-plugin"
    build_package(project, package)
    monkeypatch.setenv("ATLAS_ALLOW_UNSIGNED_PLUGINS", "true")

    result = invoke(
        TestClient(create_app()), external_payload(package.read_bytes())
    ).json()

    assert result == {
        "ok": False,
        "error": {"code": "runner_output_schema_invalid"},
    }


def test_external_plugin_cannot_import_undeclared_runner_dependency(
    tmp_path, monkeypatch
):
    package = unsigned_package(tmp_path, monkeypatch, '''
import cryptography

class Plugin:
    async def parse(self, request, context):
        if False:
            yield None
''')

    result = invoke(TestClient(create_app()), external_payload(package)).json()

    assert result == {
        "ok": False,
        "error": {"code": "plugin_dependency_not_declared", "type": "ImportError"},
    }


def test_external_plugin_cannot_spoof_filename_to_import_undeclared_dependency(
    tmp_path, monkeypatch
):
    package = unsigned_package(tmp_path, monkeypatch, '''
exec(compile("import cryptography", "trusted_generated.py", "exec"))

class Plugin:
    async def parse(self, request, context):
        if False:
            yield None
''')

    result = invoke(TestClient(create_app()), external_payload(package)).json()

    assert result == {
        "ok": False,
        "error": {"code": "plugin_dependency_not_declared", "type": "ImportError"},
    }


def test_external_plugin_cannot_mutate_module_metadata_to_hide_provenance(
    tmp_path, monkeypatch
):
    package = unsigned_package(tmp_path, monkeypatch, '''
__file__ = "trusted_generated.py"
__spec__ = None
import cryptography

class Plugin:
    async def parse(self, request, context):
        if False:
            yield None
''')

    result = invoke(TestClient(create_app()), external_payload(package)).json()

    assert result == {
        "ok": False,
        "error": {"code": "plugin_dependency_not_declared", "type": "ImportError"},
    }


def test_external_plugin_cannot_exec_import_with_trusted_globals(
    tmp_path, monkeypatch
):
    package = unsigned_package(tmp_path, monkeypatch, '''
import builtins
exec("import cryptography", vars(builtins))

class Plugin:
    async def parse(self, request, context):
        if False:
            yield None
''')

    result = invoke(TestClient(create_app()), external_payload(package)).json()

    assert result == {
        "ok": False,
        "error": {"code": "plugin_dependency_not_declared", "type": "ImportError"},
    }
